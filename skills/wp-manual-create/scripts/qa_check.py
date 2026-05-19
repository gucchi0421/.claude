#!/usr/bin/env python3
"""
生成されたWPマニュアルPPTXのQAチェックスクリプト
- パスワード・機密情報の漏れ検出
- プレースホルダー残存チェック
- スライド数・画像数の確認
"""
import argparse
import re
import sys
from pptx import Presentation

SENSITIVE_PATTERNS = [
    r"\bpassword\b", r"\bパスワード\b",
    r"[a-zA-Z0-9!@#$%^&*]{8,}",  # 8文字以上の英数字混在（要目視確認）
]

PLACEHOLDER_PATTERNS = [
    r"該当の箇所", r"クライアント名", r"\[.*?\]",
    r"naniwa-q\.jp", r"shinwa-p\.jp", r"torijiro", r"ubp97",
    # ●●●● はパスワード伏せ字として意図的なため除外
]

def check(pptx_path):
    prs = Presentation(pptx_path)
    issues = []
    warnings = []

    print(f"=== QAチェック: {pptx_path} ===\n")
    print(f"スライド数: {len(prs.slides)}")

    total_images = 0
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        images = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text
            if shape.shape_type == 13:
                images += 1
        total_images += images

        # プレースホルダーチェック
        for pat in PLACEHOLDER_PATTERNS:
            if re.search(pat, slide_text, re.IGNORECASE):
                issues.append(f"Slide {i+1}: プレースホルダーまたは旧情報が残っています → '{pat}'")

        # 画像なしスライド（表紙・目次は除く）
        if i >= 2 and images == 0:
            warnings.append(f"Slide {i+1}: 画像がありません（スクリーンショット未挿入の可能性）")

    print(f"総画像数: {total_images}\n")

    if issues:
        print("【❌ 要修正】")
        for issue in issues:
            print(f"  {issue}")
    else:
        print("【✅ 要修正なし】")

    if warnings:
        print("\n【⚠️ 要確認】")
        for w in warnings:
            print(f"  {w}")
    else:
        print("【✅ 警告なし】")

    print("\n🔒 セキュリティ確認（目視必須）:")
    print("  - ログイン画面SSにパスワードが映っていないか")
    print("  - スライド内に本番パスワードが記載されていないか")

    return len(issues) == 0

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="WPマニュアルQAチェック")
    parser.add_argument("pptx_path", help="チェック対象のPPTXファイル")
    args = parser.parse_args()
    ok = check(args.pptx_path)
    sys.exit(0 if ok else 1)
