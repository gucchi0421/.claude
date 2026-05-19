#!/usr/bin/env python3
"""
wp-manual-create: PPTXマニュアル生成スクリプト
雛形PPTXにスクリーンショットとテキストを埋め込んで完成マニュアルを生成する
"""

import argparse
import json
import os
import copy
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import lxml.etree as etree


# ============================================================
# スライド役割の自動判定
# ============================================================

SLIDE_ROLE_RULES = [
    ("TOC_SLIDE",       ["目次", "contents", "もくじ"]),
    ("LOGIN_SLIDE",     ["ログイン", "login", "wp-admin"]),
    ("POST_LIST_SLIDE", ["投稿", "一覧", "新規追加", "編集", "ゴミ箱"]),
    ("POST_EDIT_SLIDE", ["タイトル", "本文", "公開", "ビジュアル", "プレビュー"]),
    ("MEDIA_SLIDE",     ["メディア", "画像", "アイキャッチ", "アップロード"]),
    ("EDITOR_SLIDE",    ["エディター", "書式", "太字", "テキスト編集"]),
    ("CUSTOM_POST",     []),  # カスタム投稿用（動的に割り当て）
]

def detect_slide_role(slide):
    """スライドのテキストから役割を判定する"""
    all_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            all_text += shape.text_frame.text.lower()

    for role, keywords in SLIDE_ROLE_RULES:
        if role == "CUSTOM_POST":
            continue
        for kw in keywords:
            if kw.lower() in all_text:
                return role
    return "GENERIC_SLIDE"


# ============================================================
# テキスト置換
# ============================================================

REPLACEMENTS_MAP = {
    # ログインスライドのURL置換対象パターン
    "url_placeholder": [
        r"https?://[^\s/]+/wp(-admin)?/?",
        r"naniwa-q\.jp",
        r"shinwa-p\.jp",
        r"example\.com",
    ],
    # ユーザー名/パスワード行の置換（パスワードは伏せる）
    "credential_line": [
        r"ユーザー名[：:]\s*\S+",
        r"パスワード\s*[：:]\s*\S+",
        r"user(name)?[：:]\s*\S+",
        r"pass(word)?[：:]\s*\S+",
    ],
}

def replace_text_in_shape(shape, old_pattern, new_text):
    """シェイプ内のテキストをパターンマッチで置換する（書式を保持）"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            new = re.sub(old_pattern, new_text, run.text)
            if new != run.text:
                run.text = new

def update_login_slide(slide, admin_url):
    """ログインスライドのURLを更新し、パスワードを伏せる"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # URLを置換
        for pattern in REPLACEMENTS_MAP["url_placeholder"]:
            replace_text_in_shape(shape, pattern, admin_url.rstrip('/'))
        # パスワードを伏せる
        for pattern in REPLACEMENTS_MAP["credential_line"]:
            # ユーザー名はそのまま、パスワードは●●●●●●
            if "パスワード" in shape.text_frame.text or "pass" in shape.text_frame.text.lower():
                replace_text_in_shape(shape, r"(パスワード\s*[：:]\s*)\S+", r"\1●●●●●●")
            if "ユーザー名" in shape.text_frame.text or "user" in shape.text_frame.text.lower():
                # ユーザー名はクライアントに確認してもらうため [ユーザー名] とする
                replace_text_in_shape(shape, r"(ユーザー名[：:]\s*)\S+", r"\1[ユーザー名]")

def update_toc_slide(slide, custom_posts, base_page_num=4):
    """目次スライドにカスタム投稿タイプを追加する"""
    if not custom_posts:
        return
    # 末尾のテキストボックスを探して追加行を挿入
    for shape in slide.shapes:
        if shape.has_text_frame and "ログイン" in shape.text_frame.text:
            tf = shape.text_frame
            # カスタム投稿タイプの目次行を追加
            for i, cp in enumerate(custom_posts):
                p = tf.add_paragraph()
                p.text = f"{len(tf.paragraphs)}.　{cp['name']}の更新方法・・・P{base_page_num + i * 2}"
                if tf.paragraphs[1].runs:
                    p.runs[0].font.size = tf.paragraphs[1].runs[0].font.size
            break


# ============================================================
# スクリーンショット挿入
# ============================================================

def find_main_picture(slide):
    """スライド内の最も大きい画像シェイプを返す"""
    pictures = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    if not pictures:
        return None
    return max(pictures, key=lambda s: s.width * s.height)

def replace_screenshot(slide, image_path):
    """スライドの主要画像をスクリーンショットで置き換える"""
    if not os.path.exists(image_path):
        print(f"  ⚠️  画像が見つかりません: {image_path}")
        return False

    target = find_main_picture(slide)
    if not target:
        # 画像シェイプがない場合は新規追加（スライド中央に配置）
        slide.shapes.add_picture(
            image_path,
            Inches(0.5), Inches(1.8),
            Inches(9.0), Inches(5.0)
        )
        return True

    # 既存画像と同じ位置・サイズで差し替え
    left, top, width, height = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(image_path, left, top, width, height)
    return True


# ============================================================
# カスタム投稿スライドの複製・生成
# ============================================================

def duplicate_slide(prs, slide_index):
    """指定インデックスのスライドを複製して末尾に追加する"""
    template = prs.slides[slide_index]
    blank_layout = template.slide_layout

    # XMLレベルでスライドを複製
    xml_str = template._element.xml
    new_slide = prs.slides.add_slide(blank_layout)

    # 既存シェイプを削除して複製内容で上書き
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        if child.tag.endswith('}sp') or child.tag.endswith('}pic') or child.tag.endswith('}graphicFrame'):
            sp_tree.remove(child)

    # 元スライドのシェイプを複製
    orig_sp_tree = template.shapes._spTree
    for child in orig_sp_tree:
        if not child.tag.endswith('}sp') and not child.tag.endswith('}pic'):
            continue
        new_sp_tree_elem = copy.deepcopy(child)
        sp_tree.append(new_sp_tree_elem)

    return new_slide

def create_custom_post_slides(prs, custom_post, list_slide_idx, edit_slide_idx, screenshots_dir):
    """カスタム投稿タイプ用のスライドペア（一覧+編集）を生成する"""
    print(f"  カスタム投稿スライド生成: {custom_post['name']}")

    # --- 一覧スライド ---
    list_slide = duplicate_slide(prs, list_slide_idx)
    for shape in list_slide.shapes:
        if shape.has_text_frame:
            # セクションタイトルを更新
            txt = shape.text_frame.text
            if "更新情報" in txt or "投稿" in txt or "NEWS" in txt:
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = f"{custom_post['name']}の更新方法 -1"
                break

    list_ss = os.path.join(screenshots_dir, custom_post['screenshots'].get('list', ''))
    if list_ss and os.path.exists(list_ss):
        replace_screenshot(list_slide, list_ss)

    # --- 編集スライド ---
    edit_slide = duplicate_slide(prs, edit_slide_idx)
    for shape in edit_slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if "更新情報" in txt or "投稿" in txt or "NEWS" in txt:
                for run in shape.text_frame.paragraphs[0].runs:
                    run.text = f"{custom_post['name']}の更新方法 -2"
                break

    edit_ss = os.path.join(screenshots_dir, custom_post['screenshots'].get('edit', ''))
    if edit_ss and os.path.exists(edit_ss):
        replace_screenshot(edit_slide, edit_ss)

        # カスタムフィールドがある場合は注記を追加
        if custom_post.get('has_custom_fields'):
            fields_text = "・カスタム項目: " + "、".join(custom_post.get('fields', []))
            for shape in edit_slide.shapes:
                if shape.has_text_frame and len(shape.text_frame.text) > 20:
                    p = shape.text_frame.add_paragraph()
                    p.text = fields_text
                    break

    return list_slide, edit_slide


# ============================================================
# メイン処理
# ============================================================

def generate_manual(template_path, site_info_path, screenshots_dir, output_path):
    print("=== WP更新マニュアル生成開始 ===\n")

    # サイト情報読み込み
    with open(site_info_path) as f:
        site_info = json.load(f)

    site_name = site_info.get('site_name', 'クライアント')
    admin_url = site_info.get('admin_url', '')
    custom_posts = site_info.get('custom_posts', [])

    print(f"サイト名: {site_name}")
    print(f"カスタム投稿タイプ: {[cp['name'] for cp in custom_posts]}\n")

    # テンプレート読み込み
    prs = Presentation(template_path)

    # スライドの役割を判定
    slide_roles = {}
    for i, slide in enumerate(prs.slides):
        role = detect_slide_role(slide)
        slide_roles[i] = role
        print(f"  Slide {i+1}: {role}")

    print()

    # 役割インデックスのマップ
    role_idx = {}
    for idx, role in slide_roles.items():
        role_idx.setdefault(role, []).append(idx)

    # --- 表紙スライド ---
    title_slides = role_idx.get('GENERIC_SLIDE', [])
    if title_slides:
        title_slide = prs.slides[title_slides[0]]
        for shape in title_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = re.sub(r'汎用（クライアント名）|クライアント名|（.*?）', site_name, run.text)

    # --- ログインスライド ---
    for idx in role_idx.get('LOGIN_SLIDE', []):
        print(f"ログインスライド更新 (Slide {idx+1})")
        update_login_slide(prs.slides[idx], admin_url)
        ss_path = os.path.join(screenshots_dir, 'ss_01_dashboard.png')
        replace_screenshot(prs.slides[idx], ss_path)

    # --- 目次スライド ---
    for idx in role_idx.get('TOC_SLIDE', []):
        print(f"目次スライド更新 (Slide {idx+1})")
        update_toc_slide(prs.slides[idx], custom_posts)

    # --- 投稿一覧スライド ---
    post_list_indices = role_idx.get('POST_LIST_SLIDE', [])
    for idx in post_list_indices:
        ss_path = os.path.join(screenshots_dir, 'ss_02_post_list.png')
        replace_screenshot(prs.slides[idx], ss_path)

    # --- 投稿編集スライド ---
    post_edit_indices = role_idx.get('POST_EDIT_SLIDE', [])
    for idx in post_edit_indices:
        ss_path = os.path.join(screenshots_dir, 'ss_03_post_edit.png')
        replace_screenshot(prs.slides[idx], ss_path)

    # --- メディアスライド ---
    for idx in role_idx.get('MEDIA_SLIDE', []):
        ss_path = os.path.join(screenshots_dir, 'ss_05_media_upload.png')
        replace_screenshot(prs.slides[idx], ss_path)

    # --- カスタム投稿スライド追加 ---
    if custom_posts and post_list_indices and post_edit_indices:
        list_tmpl_idx = post_list_indices[-1]
        edit_tmpl_idx = post_edit_indices[-1]

        for cp in custom_posts:
            create_custom_post_slides(
                prs, cp, list_tmpl_idx, edit_tmpl_idx, screenshots_dir
            )

    # --- 出力 ---
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"\n✅ マニュアル生成完了: {output_path}")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="WP更新マニュアル生成")
    parser.add_argument("--template", required=True, help="雛形PPTXのパス")
    parser.add_argument("--site-info", required=True, help="サイト情報JSONのパス")
    parser.add_argument("--screenshots-dir", required=True, help="スクリーンショット保存ディレクトリ")
    parser.add_argument("--output", required=True, help="出力PPTXのパス")
    args = parser.parse_args()

    generate_manual(
        args.template,
        args.site_info,
        args.screenshots_dir,
        args.output
    )
