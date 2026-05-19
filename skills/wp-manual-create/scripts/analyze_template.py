#!/usr/bin/env python3
"""
雛形PPTXの構造を解析してtemplate_map.jsonを出力する
"""
import argparse, json, re, os
from pptx import Presentation

ROLE_RULES = [
    ("TOC_SLIDE",       ["目次", "contents"]),
    ("LOGIN_SLIDE",     ["ログイン", "login", "wp-admin", "ユーザー名", "パスワード"]),
    ("POST_LIST_SLIDE", ["投稿", "一覧", "新規追加", "ゴミ箱", "編集"]),
    ("POST_EDIT_SLIDE", ["タイトル", "本文", "公開", "ビジュアル", "プレビュー", "下書き"]),
    ("MEDIA_SLIDE",     ["メディア", "画像", "アイキャッチ", "アップロード", "ドラッグ"]),
    ("EDITOR_SLIDE",    ["エディター", "書式", "太字", "テキスト編集", "斜体", "リンク"]),
]

def detect_role(slide):
    text = " ".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    ).lower()
    for role, keywords in ROLE_RULES:
        if any(k.lower() in text for k in keywords):
            return role
    return "GENERIC_SLIDE"

def analyze(pptx_path):
    prs = Presentation(pptx_path)
    result = {
        "template_path": pptx_path,
        "slide_count": len(prs.slides),
        "slide_size": {
            "width_inches": round(prs.slide_width.inches, 2),
            "height_inches": round(prs.slide_height.inches, 2),
            "orientation": "landscape" if prs.slide_width > prs.slide_height else "portrait"
        },
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        role = detect_role(slide)
        texts = [shape.text_frame.text[:80] for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
        pictures = [{"name": s.name, "w": round(s.width.inches,2), "h": round(s.height.inches,2)}
                    for s in slide.shapes if s.shape_type == 13]
        result["slides"].append({
            "index": i,
            "slide_number": i + 1,
            "role": role,
            "text_preview": texts[:3],
            "pictures": pictures,
            "picture_count": len(pictures)
        })

    return result

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx_path")
    parser.add_argument("--output", default="template_map.json")
    args = parser.parse_args()

    result = analyze(args.pptx_path)
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"解析完了: {args.output}")
    print(f"スライド数: {result['slide_count']}")
    print(f"サイズ: {result['slide_size']['width_inches']}\" x {result['slide_size']['height_inches']}\" ({result['slide_size']['orientation']})")
    print()
    for s in result["slides"]:
        print(f"  Slide {s['slide_number']:2d}: {s['role']:20s} | 画像{s['picture_count']}枚 | {s['text_preview'][0][:40] if s['text_preview'] else '(テキストなし)'}")
